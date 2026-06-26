import os
import shutil
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether, HRFlowable
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen import canvas

import docx
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import nsdecls, qn

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Define output directories
DOCS_DIR = r"docs"
os.makedirs(DOCS_DIR, exist_ok=True)

# Colors definition
PRIMARY_HEX = "#0078D4"  # Azure Blue
SECONDARY_HEX = "#005A9E"  # Deep Blue
ACCENT_HEX = "#00B7C3"  # Teal / Cyan
DARK_NEUTRAL_HEX = "#0F172A"  # Slate-900
LIGHT_NEUTRAL_HEX = "#F8FAFC"  # Slate-50
BORDER_COLOR_HEX = "#CBD5E1"  # Slate-300
TEXT_MUTED_HEX = "#64748B"  # Slate-500
SUCCESS_HEX = "#107C41"  # Excel Green
CALLOUT_BG_HEX = "#F0F7FF"  # Very Light Azure

# ---------------------------------------------------------
# PDF GENERATION USING REPORTLAB
# ---------------------------------------------------------

class NumberedCanvas(canvas.Canvas):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self._saved_page_states = []

    def showPage(self):
        self._saved_page_states.append(dict(self.__dict__))
        self._startPage()

    def save(self):
        num_pages = len(self._saved_page_states)
        for state in self._saved_page_states:
            self.__dict__.update(state)
            self.draw_page_elements(num_pages)
            super().showPage()
        super().save()

    def draw_page_elements(self, page_count):
        # Skip cover page
        if self._pageNumber == 1:
            return
        
        self.saveState()
        self.setFont("Helvetica", 8)
        self.setFillColor(colors.HexColor(TEXT_MUTED_HEX))
        
        # Header
        self.drawString(54, 750, "AnalytiX Platform - Inputs & Outputs Reference Guide")
        self.setStrokeColor(colors.HexColor(BORDER_COLOR_HEX))
        self.setLineWidth(0.5)
        self.line(54, 742, 558, 742)
        
        # Footer
        self.line(54, 45, 558, 45)
        self.drawString(54, 30, "CONFIDENTIAL - For AnalytiX Enterprise Evaluation Only")
        page_text = f"Page {self._pageNumber} of {page_count}"
        self.drawRightString(558, 30, page_text)
        
        self.restoreState()

def draw_pdf_cover_background(canvas_obj, doc):
    canvas_obj.saveState()
    # Dark blue background banner (top 45% of cover page)
    canvas_obj.setFillColor(colors.HexColor(DARK_NEUTRAL_HEX))
    canvas_obj.rect(0, 420, 612, 372, stroke=0, fill=1)
    
    # Bottom brand accent line (Azure blue)
    canvas_obj.setStrokeColor(colors.HexColor(PRIMARY_HEX))
    canvas_obj.setLineWidth(5)
    canvas_obj.line(54, 110, 558, 110)
    
    canvas_obj.restoreState()

def create_pdf_styles():
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'CoverTitle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=28,
        leading=34,
        textColor=colors.HexColor('#ffffff'),
        spaceAfter=12
    )
    
    subtitle_style = ParagraphStyle(
        'CoverSubtitle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=12,
        leading=16,
        textColor=colors.HexColor(ACCENT_HEX),
        spaceAfter=15
    )
    
    h1_style = ParagraphStyle(
        'SectionH1',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=22,
        textColor=colors.HexColor(PRIMARY_HEX),
        spaceBefore=18,
        spaceAfter=8,
        keepWithNext=True
    )
    
    h2_style = ParagraphStyle(
        'SectionH2',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=13,
        leading=16,
        textColor=colors.HexColor(SECONDARY_HEX),
        spaceBefore=12,
        spaceAfter=6,
        keepWithNext=True
    )
    
    body_style = ParagraphStyle(
        'Body',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor('#334155'), # slate-700
        spaceBefore=4,
        spaceAfter=6
    )
    
    bullet_style = ParagraphStyle(
        'Bullet',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor('#334155'),
        leftIndent=15,
        firstLineIndent=-10,
        spaceBefore=2,
        spaceAfter=2
    )
    
    callout_style = ParagraphStyle(
        'Callout',
        parent=styles['Normal'],
        fontName='Helvetica-Oblique',
        fontSize=9.5,
        leading=13.5,
        textColor=colors.HexColor('#1E293B'),
        backColor=colors.HexColor(CALLOUT_BG_HEX),
        borderColor=colors.HexColor(PRIMARY_HEX),
        borderWidth=0.5,
        borderPadding=10,
        spaceBefore=8,
        spaceAfter=8
    )

    journey_style = ParagraphStyle(
        'JourneyBox',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9.5,
        leading=14,
        textColor=colors.HexColor('#1E3A8A'), # Dark Blue text
        backColor=colors.HexColor('#EFF6FF'), # Soft blue background
        borderColor=colors.HexColor('#3B82F6'), # Blue border
        borderWidth=0.5,
        borderPadding=10,
        spaceBefore=8,
        spaceAfter=8
    )
    
    return {
        'title': title_style,
        'subtitle': subtitle_style,
        'h1': h1_style,
        'h2': h2_style,
        'body': body_style,
        'bullet': bullet_style,
        'callout': callout_style,
        'journey': journey_style
    }

def generate_pdf_document(filepath):
    print(f"Generating PDF: {filepath}")
    doc = SimpleDocTemplate(
        filepath,
        pagesize=letter,
        leftMargin=54,
        rightMargin=54,
        topMargin=54,
        bottomMargin=54
    )
    
    styles = create_pdf_styles()
    story = []
    
    # ------------------ COVER PAGE ------------------
    story.append(Spacer(1, 140))
    story.append(Paragraph("ANALYTIX PLATFORM", styles['title']))
    story.append(Paragraph("Enterprise Inputs & Outputs Reference Guide", styles['subtitle']))
    story.append(Spacer(1, 160))
    
    meta_text = f"""
    <b>Document Type:</b> Product Technical Overview<br/>
    <b>Classification:</b> CONFIDENTIAL - Internal & Client Review Only<br/>
    <b>Target Environment:</b> Microsoft Azure Marketplace Integration<br/>
    <b>Product Version:</b> v4.2.0 (Production-Ready Release)<br/>
    <b>Date:</b> June 2026<br/>
    <b>Publisher:</b> AnalytiX Product Marketing & Compliance Division
    """
    story.append(Paragraph(meta_text, styles['body']))
    story.append(PageBreak())
    
    # ------------------ SECTION 1 ------------------
    story.append(Paragraph("Section 1: Customer Data Inputs", styles['h1']))
    intro_p1 = """AnalytiX integrates directly into existing life-sciences IT landscapes. By leveraging advanced data connector adapters and schemas, the platform acts as a unified ingestion fabric. It seamlessly handles structural chemistry, high-throughput biology, enterprise databases, and unstructured scientific files."""
    story.append(Paragraph(intro_p1, styles['body']))
    story.append(Spacer(1, 8))
    
    # Customer Inputs Table
    input_data = [
        ["Category", "Supported Sources & Protocols", "Ingested Data Types"],
        [
            "Laboratory Systems",
            "• ELN: Benchling, Dotmatics, Signals Notebook\n• LIMS: LabWare, LabVantage, STARLIMS",
            "Experimental protocols, sample runs, inventory mappings, user audit logs."
        ],
        [
            "Enterprise Databases",
            "• Relational & Warehouses: PostgreSQL, Oracle, SQL Server, Snowflake, MySQL",
            "Federated query views, historical screening data, assay metadata databases."
        ],
        [
            "Scientific Files",
            "• File Ingest: CSV, Excel, JSON, XML, FASTA",
            "Raw assay values, plate layouts, genomic/nucleotide data, physical descriptions."
        ],
        [
            "Enterprise Integration",
            "• Connectivity: REST APIs, SOAP APIs, FTP/SFTP, Cloud Storage (AWS S3, Azure Blob, GCS)",
            "Robot file drops, vendor data deliveries, automated system-to-system syncs."
        ],
        [
            "Scientific Data",
            "• Biological & Chemical Entities: SMILES, SDF, Sequence files (DNA, RNA, Protein), Assay outputs",
            "Chemical structures, sequence alignments, inhibition values, metadata tags, audit ledgers."
        ]
    ]
    
    # Format table for PDF
    formatted_table_data = []
    # Header
    formatted_table_data.append([Paragraph(cell, ParagraphStyle('TH', parent=styles['body'], fontName='Helvetica-Bold', textColor=colors.white)) for cell in input_data[0]])
    # Data Rows
    for row in input_data[1:]:
        formatted_table_data.append([
            Paragraph(row[0], ParagraphStyle('TC_B', parent=styles['body'], fontName='Helvetica-Bold')),
            Paragraph(row[1].replace("\n", "<br/>"), styles['body']),
            Paragraph(row[2], styles['body'])
        ])
        
    input_table = Table(formatted_table_data, colWidths=[100, 200, 204])
    input_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor(PRIMARY_HEX)),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,1), (-1,-1), 6),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor(BORDER_COLOR_HEX)),
        ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#F8FAFC')),
    ]))
    story.append(input_table)
    story.append(Spacer(1, 15))
    
    # ------------------ SECTION 2 ------------------
    story.append(Paragraph("Section 2: Platform Data Processing", styles['h1']))
    story.append(Paragraph("Upon data ingestion, the AnalytiX platform executes sequential pipelines to normalize, register, link, and analyze scientific entities in a regulated, high-performance environment.", styles['body']))
    
    # Processing list
    processing_steps = [
        "<b>Data Ingestion & Validation:</b> Performs type verification, validation rules checks, and structural normalization on compounds and sequences.",
        "<b>Schema & Metadata Discovery:</b> Automatically inspects connected databases to map structures and index attributes into a dynamic EAV Metadata Catalog.",
        "<b>Federated Query Execution:</b> Translates queries across distributed databases (e.g. Snowflake + local PostgreSQL) without data replication.",
        "<b>Scientific Analysis:</b> Runs RDKit chemical descriptors calculations, global/local sequence alignments, and dose-response curve regressions (4PL models).",
        "<b>Workflow Automation & AI:</b> Executes multi-stage pipelines (e.g., query candidate -> filter ADMET -> dock molecule -> request compliance signature) with AI Copilot support.",
        "<b>Lineage & Auditing:</b> Maintains complete end-to-end data lineage from source data connector to final fit, writing to a cryptographically linked SHA-256 ledger."
    ]
    for step in processing_steps:
        story.append(Paragraph(f"• {step}", styles['bullet']))
        
    story.append(Spacer(1, 10))
    story.append(Paragraph("Platform Processing Architecture Flow", styles['h2']))
    
    # Flow diagram as table
    flow_steps = ["Raw Input\nData", "Connector\nLayer", "Metadata\nCatalog", "Federated\nQuery Engine", "Scientific\nAnalytics", "AI Scientist\nCopilot", "Reports &\nDashboards"]
    flow_table_data = [[Paragraph(step.replace("\n", "<br/>"), ParagraphStyle('FS', parent=styles['body'], fontName='Helvetica-Bold', fontSize=8, leading=10, alignment=1, textColor=colors.HexColor('#0F172A'))) for step in flow_steps]]
    
    flow_table = Table(flow_table_data, colWidths=[72]*7)
    flow_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,0), colors.HexColor('#E2E8F0')),
        ('BACKGROUND', (1,0), (1,0), colors.HexColor('#CBD5E1')),
        ('BACKGROUND', (2,0), (2,0), colors.HexColor('#94A3B8')),
        ('BACKGROUND', (3,0), (3,0), colors.HexColor('#64748B')),
        ('BACKGROUND', (4,0), (4,0), colors.HexColor('#3B82F6')), # Blue accent
        ('BACKGROUND', (5,0), (5,0), colors.HexColor('#1D4ED8')), # Deep Blue
        ('BACKGROUND', (6,0), (6,0), colors.HexColor(SUCCESS_HEX)), # Green accent
        ('TEXTCOLOR', (4,0), (6,0), colors.white),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('GRID', (0,0), (-1,-1), 1.5, colors.HexColor('#FFFFFF')),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
    ]))
    
    # Update text styles for flow table (white text on dark backgrounds)
    for col_idx in [4, 5, 6]:
        flow_table_data[0][col_idx].style.textColor = colors.white
        
    story.append(flow_table)
    story.append(Spacer(1, 15))
    
    # ------------------ SECTION 3 ------------------
    story.append(Paragraph("Section 3: Platform Data Outputs", styles['h1']))
    story.append(Paragraph("AnalytiX outputs are tailored to support scientific discoveries, project leadership review, and strict regulatory inspections.", styles['body']))
    
    outputs_data = [
        ("Scientific Analytics", "Dose-response sigmoidal curves (IC50/EC50 fits), molecular similarity matrices (Tanimoto coefficients), genomic sequence identity grids, R-group distributions, and correlation maps."),
        ("AI Copilot Outputs", "Natural language answers to databases queries, automated study summaries, compound optimization recommendations, and workflow template generations."),
        ("Enterprise Reports", "High-fidelity PDF compliance files, structured Excel summaries (with raw assay metrics), PPTX presentations for project review, and regulatory audit-ready logs."),
        ("Compliance Records", "FDA 21 CFR Part 11 electronic signature validations, SHA-256 cryptographic audit logs, and interactive interactive visual data lineage tracking reports."),
        ("Operational Alerts", "Workflow run status notices, automated approval request notifications, and scheduled pipeline reports delivered via Webhooks, Email, or Slack.")
    ]
    
    for title, desc in outputs_data:
        story.append(Paragraph(f"<b>{title}:</b> {desc}", styles['body']))
        
    story.append(Spacer(1, 15))
    story.append(PageBreak())
    
    # ------------------ SECTION 4 ------------------
    story.append(Paragraph("Section 4: Enterprise Business Benefits", styles['h1']))
    story.append(Paragraph("The platform maximizes value across all departments, establishing a modern, compliant, and data-driven discovery pipeline.", styles['body']))
    
    # Benefits Grid Table
    benefits_data = [
        [
            Paragraph("<b>RESEARCH SCIENTISTS</b><br/><i>Accelerated In Silico Discovery</i><br/>• Saves hours of database manual querying and scripting.<br/>• Correlates chemical structures with biology assays automatically.<br/>• Natural language interface allows complex queries without SQL.", styles['callout']),
            Paragraph("<b>RESEARCH MANAGEMENT</b><br/><i>Data-Driven Portfolio Direction</i><br/>• Real-time project dashboard trackers and KPIs.<br/>• Fast candidate optimization cycle analysis.<br/>• Centralized decision support for candidate selection.", styles['callout'])
        ],
        [
            Paragraph("<b>COMPLIANCE TEAMS</b><br/><i>Audit-Ready Transparency</i><br/>• Strict FDA 21 CFR Part 11 compliant digital logs.<br/>• Complete lineage mappings from raw source to final charts.<br/>• Immutable cryptographic audit trailing blocks.", styles['callout']),
            Paragraph("<b>IT & INFRASTRUCTURE</b><br/><i>Secure, Integrated Ecosystem</i><br/>• Connects directly to Snowflake, PostgreSQL, ELN, and LIMS.<br/>• Granular RBAC, SSO, and OAuth standard integrations.<br/>• High availability containerized microservices structure.", styles['callout'])
        ]
    ]
    
    # Override background colors for different benefits blocks
    # Set the style properties inside paragraph elements
    benefits_table = Table(benefits_data, colWidths=[248, 248])
    benefits_table.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
        ('RIGHTPADDING', (0,0), (-1,-1), 8),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
    ]))
    story.append(benefits_table)
    story.append(Spacer(1, 15))
    
    # ------------------ SECTION 5 ------------------
    story.append(Paragraph("Section 5: Sample Client Journey", styles['h1']))
    journey_p = """<b>Scenario:</b> A biopharma client connects their PostgreSQL screening database, Benchling ELN protocols, cell-line CSV results, and a set of FASTA sequences.<br/><br/>
    <b>AnalytiX Automated Sequence:</b><br/>
    1. <b>Connects:</b> Securely establishes read-only data connections in minutes.<br/>
    2. <b>Discovers:</b> Auto-maps database tables and extracts custom metadata fields.<br/>
    3. <b>Relates:</b> Link assay data plates directly to compound chemical profiles using SMILES mappings.<br/>
    4. <b>Analyzes:</b> Computes dose-response parameters (IC50 curves) and registers chemical structural descriptors.<br/>
    5. <b>AI Grounds:</b> Scientist Copilot translates natural language questions into secure queries, returning analog candidate structures and compliance PDF summaries ready for review."""
    story.append(Paragraph(journey_p, styles['journey']))
    story.append(Spacer(1, 15))
    
    # ------------------ SECTION 6 ------------------
    story.append(Paragraph("Section 6: Visual Infographic Summary", styles['h1']))
    story.append(Paragraph("The AnalytiX value chain transforms customer data investments into validated business and scientific outcomes.", styles['body']))
    
    # Infographic table
    info_data = [
        ["CUSTOMER DATA", "➔", "AnalytiX PLATFORM", "➔", "SCIENTIFIC INTELLIGENCE", "➔", "BUSINESS OUTCOMES"],
        [
            "ELN & LIMS Systems\nDatabases & Files\nREST/SOAP APIs\nRaw Structures",
            "",
            "Schema Discovery\nData Standardization\nFederated Engine\nData Lineage Map",
            "",
            "AI Scientist Copilot\n4PL Dose-Response\nSequence Alignments\nCompound Analytics",
            "",
            "Faster Lead ID\n100% Compliance\nData Modernization\nLow IT Silos"
        ]
    ]
    
    info_formatted = []
    # Title Row
    info_formatted.append([
        Paragraph(cell, ParagraphStyle('IT', parent=styles['body'], fontName='Helvetica-Bold', fontSize=9, alignment=1, textColor=colors.white))
        if cell != "➔" else Paragraph(cell, ParagraphStyle('IA', parent=styles['body'], fontName='Helvetica-Bold', fontSize=12, alignment=1, textColor=colors.HexColor(PRIMARY_HEX)))
        for cell in info_data[0]
    ])
    # Content Row
    info_formatted.append([
        Paragraph(cell.replace("\n", "<br/>"), ParagraphStyle('IC', parent=styles['body'], fontSize=8.5, leading=12, alignment=1))
        if cell != "" else Paragraph("", styles['body'])
        for cell in info_data[1]
    ])
    
    info_table = Table(info_formatted, colWidths=[105, 20, 105, 20, 115, 20, 115])
    info_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (0,0), colors.HexColor(DARK_NEUTRAL_HEX)),
        ('BACKGROUND', (2,0), (2,0), colors.HexColor(SECONDARY_HEX)),
        ('BACKGROUND', (4,0), (4,0), colors.HexColor(PRIMARY_HEX)),
        ('BACKGROUND', (6,0), (6,0), colors.HexColor(SUCCESS_HEX)),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('GRID', (0,0), (0,1), 0.5, colors.HexColor(BORDER_COLOR_HEX)),
        ('GRID', (2,0), (2,1), 0.5, colors.HexColor(BORDER_COLOR_HEX)),
        ('GRID', (4,0), (4,1), 0.5, colors.HexColor(BORDER_COLOR_HEX)),
        ('GRID', (6,0), (6,1), 0.5, colors.HexColor(BORDER_COLOR_HEX)),
        ('BACKGROUND', (0,1), (0,1), colors.HexColor('#F8FAFC')),
        ('BACKGROUND', (2,1), (2,1), colors.HexColor('#F0F7FF')),
        ('BACKGROUND', (4,1), (4,1), colors.HexColor('#EFF6FF')),
        ('BACKGROUND', (6,1), (6,1), colors.HexColor('#ECFDF5')),
        ('TOPPADDING', (0,0), (-1,-1), 8),
        ('BOTTOMPADDING', (0,0), (-1,-1), 8),
    ]))
    
    story.append(info_table)
    
    # Build document
    doc.build(story, onFirstPage=draw_pdf_cover_background, canvasmaker=NumberedCanvas)
    print("PDF generated successfully.")

# ---------------------------------------------------------
# DOCX GENERATION USING PYTHON-DOCX
# ---------------------------------------------------------

def set_cell_bg(cell, hex_color):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{hex_color}"/>')
    tcPr.append(shd)

def set_callout_borders(cell, border_color_hex="0078D4"):
    tcPr = cell._tc.get_or_add_tcPr()
    tcBorders = parse_xml(f'''
        <w:tcBorders {nsdecls("w")}>
            <w:top w:val="none"/>
            <w:left w:val="single" w:sz="36" w:space="0" w:color="{border_color_hex}"/>
            <w:bottom w:val="none"/>
            <w:right w:val="none"/>
        </w:tcBorders>
    ''')
    tcPr.append(tcBorders)

def set_table_light_borders(table, border_color_hex="CBD5E1"):
    tblPr = table._tbl.tblPr
    tblBorders = parse_xml(f'''
        <w:tblBorders {nsdecls("w")}>
            <w:top w:val="single" w:sz="4" w:space="0" w:color="{border_color_hex}"/>
            <w:left w:val="none"/>
            <w:bottom w:val="single" w:sz="8" w:space="0" w:color="{border_color_hex}"/>
            <w:right w:val="none"/>
            <w:insideH w:val="single" w:sz="4" w:space="0" w:color="{border_color_hex}"/>
            <w:insideV w:val="none"/>
        </w:tblBorders>
    ''')
    tblPr.append(tblBorders)

def add_docx_heading(doc, text, level, space_before=12, space_after=6):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(space_before)
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.keep_with_next = True
    
    run = p.add_run(text)
    run.bold = True
    run.font.name = 'Segoe UI'
    if level == 1:
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0, 120, 212) # Azure Blue
    elif level == 2:
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor(0, 90, 158) # Deep Blue
    else:
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(15, 23, 42)
    return p

def generate_docx_document(filepath):
    print(f"Generating DOCX: {filepath}")
    doc = Document()
    
    # Page setup
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
        
    # Set default style properties
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Segoe UI'
    font.size = Pt(10.5)
    font.color.rgb = RGBColor(51, 65, 85) # Slate 700
    
    # ------------------ COVER PAGE ------------------
    # Add large spacing
    for _ in range(5):
        doc.add_paragraph()
        
    p_title = doc.add_paragraph()
    run_title = p_title.add_run("ANALYTIX PLATFORM")
    run_title.font.name = 'Segoe UI'
    run_title.font.size = Pt(32)
    run_title.font.bold = True
    run_title.font.color.rgb = RGBColor(11, 37, 64) # Dark slate
    
    p_sub = doc.add_paragraph()
    run_sub = p_sub.add_run("Enterprise Inputs & Outputs Reference Guide")
    run_sub.font.name = 'Segoe UI'
    run_sub.font.size = Pt(14)
    run_sub.font.color.rgb = RGBColor(0, 120, 212) # Azure blue
    p_sub.paragraph_format.space_after = Pt(24)
    
    # Metadata callout table
    meta_table = doc.add_table(rows=1, cols=1)
    meta_table.alignment = WD_TABLE_ALIGNMENT.LEFT
    cell = meta_table.cell(0, 0)
    set_cell_bg(cell, "F8FAFC")
    set_callout_borders(cell, "0078D4")
    
    meta_p = cell.paragraphs[0]
    meta_p.paragraph_format.left_indent = Inches(0.15)
    meta_p.paragraph_format.space_before = Pt(8)
    meta_p.paragraph_format.space_after = Pt(8)
    
    meta_text = (
        "Document Type: Product Technical Specification\n"
        "Classification: CONFIDENTIAL - Evaluation Only\n"
        "Azure Marketplace Reference: AnalytiX-v4.2-Overview\n"
        "Compliance Environment: FDA 21 CFR Part 11 Compliant\n"
        "Publisher: AnalytiX Product Engineering & Compliance Group\n"
        "Last Updated: June 2026"
    )
    meta_run = meta_p.add_run(meta_text)
    meta_run.font.name = 'Segoe UI'
    meta_run.font.size = Pt(9.5)
    meta_run.font.color.rgb = RGBColor(100, 116, 139)
    
    doc.add_page_break()
    
    # ------------------ SECTION 1 ------------------
    add_docx_heading(doc, "Section 1: Supported Data Inputs", level=1)
    
    doc.add_paragraph(
        "AnalytiX features modular data connector adapters designed to bridge dry-lab calculations with live systems. "
        "It acts as a secure data federation fabric, allowing scientists to access scattered data in one place."
    )
    
    inputs_table = doc.add_table(rows=6, cols=2)
    inputs_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_table_light_borders(inputs_table)
    
    headers = ["Ingestion Category", "Details & Supported Ingest Formats"]
    for i, h in enumerate(headers):
        cell = inputs_table.cell(0, i)
        set_cell_bg(cell, "0078D4")
        p = cell.paragraphs[0]
        p.paragraph_format.space_before = Pt(6)
        p.paragraph_format.space_after = Pt(6)
        run = p.add_run(h)
        run.font.name = 'Segoe UI'
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        run.font.size = Pt(10)
        
    rows_data = [
        ["Laboratory Systems", "ELN (Benchling, Dotmatics, Signals Notebook)\nLIMS (LabWare, LabVantage, STARLIMS)"],
        ["Enterprise Databases", "PostgreSQL, Oracle, SQL Server, Snowflake, MySQL"],
        ["Scientific File Formats", "CSV, Excel, JSON, XML, FASTA"],
        ["Interfaces & Cloud", "REST APIs, SOAP APIs, Cloud Storage (AWS S3, Azure Blob, GCS), Secure FTP/SFTP"],
        ["Scientific Data Objects", "Compound Structures, SMILES, Assay Results, Experimental Results, DNA/RNA/Protein sequences, Workflow Data, SHA-256 Audit Logs, User Session details"]
    ]
    
    for row_idx, data in enumerate(rows_data, 1):
        for col_idx, text in enumerate(data):
            cell = inputs_table.cell(row_idx, col_idx)
            if row_idx % 2 == 1:
                set_cell_bg(cell, "F8FAFC")
            p = cell.paragraphs[0]
            p.paragraph_format.space_before = Pt(5)
            p.paragraph_format.space_after = Pt(5)
            run = p.add_run(text)
            run.font.name = 'Segoe UI'
            run.font.size = Pt(9.5)
            if col_idx == 0:
                run.bold = True
                run.font.color.rgb = RGBColor(11, 37, 64)
                
    doc.add_paragraph().paragraph_format.space_after = Pt(10)
    
    # ------------------ SECTION 2 ------------------
    add_docx_heading(doc, "Section 2: Platform Ingest & Processing", level=1)
    
    doc.add_paragraph(
        "Ingested data is processed using containerized microservice pipelines. Tasks include metadata mapping, "
        "standardization, and computational scientific analysis in a validated FDA 21 CFR Part 11 compliant environment."
    )
    
    processing_bullets = [
        "Data Validation & Quality Checks: Ensures schema integrity and filters out malformed chemical structures.",
        "Schema & Metadata Discovery: Scans connected DBs and indexes attributes dynamically inside an EAV Catalog.",
        "Federated Query Execution: Safely queries databases across sources in real-time, eliminating local copy duplicates.",
        "Cheminformatics & Bioinformatics: Calculates molecular descriptors, global/local alignments, and 4-Parameter curve fits.",
        "Workflow Automation & Compliance: Orchestrates tasks, alerts users, logs steps to a cryptographic SHA-256 ledger."
    ]
    for b in processing_bullets:
        p = doc.add_paragraph(style='List Bullet')
        p.paragraph_format.space_after = Pt(3)
        parts = b.split(":")
        run_bold = p.add_run(parts[0] + ":")
        run_bold.bold = True
        run_bold.font.name = 'Segoe UI'
        run_normal = p.add_run(parts[1])
        run_normal.font.name = 'Segoe UI'
        
    doc.add_paragraph().paragraph_format.space_after = Pt(8)
    
    # Flow diagram
    add_docx_heading(doc, "Processing Architecture Flow:", level=2)
    p_flow = doc.add_paragraph()
    p_flow.paragraph_format.space_before = Pt(6)
    p_flow.paragraph_format.space_after = Pt(12)
    p_flow.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    flow_steps = ["Input Data", "Connector Layer", "Metadata Catalog", "Federated Query Engine", "Scientific Analytics", "AI Scientist Copilot", "Reports & Dashboards"]
    flow_str = " ➔ ".join(flow_steps)
    run_flow = p_flow.add_run(flow_str)
    run_flow.font.name = 'Segoe UI'
    run_flow.font.bold = True
    run_flow.font.size = Pt(9.5)
    run_flow.font.color.rgb = RGBColor(0, 120, 212)
    
    # ------------------ SECTION 3 ------------------
    add_docx_heading(doc, "Section 3: Platform Data Outputs", level=1)
    doc.add_paragraph(
        "AnalytiX delivers analytics, compliance reports, and operational events to key enterprise stakeholders."
    )
    
    outputs_cats = [
        ("Scientific Analytics", "Compound descriptors, 4PL curves (IC50), bio alignments, correlation heatmaps, scatter charts."),
        ("AI Scientist Outputs", "Natural language explanations, compound proposals, data summaries, automated workflows."),
        ("Enterprise Reports", "Signed PDF reports, structured Excel grids, slide decks (PPTX), and regulatory audit logs."),
        ("Regulatory Compliance", "FDA 21 CFR Part 11 records, electronic signatures validation files, and visual data lineage paths."),
        ("Operational Alerts", "Slack webhooks, email alerts, workflow execution triggers, and approval requests.")
    ]
    
    for title, desc in outputs_cats:
        p = doc.add_paragraph(style='List Bullet')
        p.paragraph_format.space_after = Pt(4)
        run_title = p.add_run(title + ": ")
        run_title.bold = True
        run_title.font.name = 'Segoe UI'
        run_title.font.color.rgb = RGBColor(0, 90, 158)
        run_desc = p.add_run(desc)
        run_desc.font.name = 'Segoe UI'
        
    doc.add_paragraph().paragraph_format.space_after = Pt(10)
    
    # ------------------ SECTION 4 ------------------
    add_docx_heading(doc, "Section 4: Enterprise Business Benefits", level=1)
    
    benefits_table = doc.add_table(rows=2, cols=2)
    benefits_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_table_light_borders(benefits_table)
    
    benefits_cells_data = [
        [
            ("RESEARCH SCIENTISTS", "• Drastically reduces manual data mining.\n• Automated structure-activity correlation.\n• AI Scientist answers natural language questions."),
            ("COMPLIANCE TEAMS", "• 100% data traceability via Data Lineage.\n• Cryptographically chain-linked audit trails.\n• Rapid generation of FDA Part 11 compliant reports.")
        ],
        [
            ("RESEARCH MANAGERS", "• Centralized key metrics and timeline dashboards.\n• Faster compound optimization screening runs.\n• Clear portfolio monitoring dashboard feeds."),
            ("IT & INFRASTRUCTURE", "• Secured SSO, OAuth, and RBAC matrix permissions.\n• Integrates with existing LIMS, ELN, and Snowflake.\n• Scalable and robust microservice deployment.")
        ]
    ]
    
    for row_idx, row in enumerate(benefits_cells_data):
        for col_idx, (title, text) in enumerate(row):
            cell = benefits_table.cell(row_idx, col_idx)
            set_cell_bg(cell, "F8FAFC")
            set_callout_borders(cell, "0078D4" if row_idx == 0 else "005A9E")
            
            p = cell.paragraphs[0]
            p.paragraph_format.left_indent = Inches(0.15)
            p.paragraph_format.space_before = Pt(6)
            p.paragraph_format.space_after = Pt(6)
            
            run_title = p.add_run(title + "\n")
            run_title.bold = True
            run_title.font.name = 'Segoe UI'
            run_title.font.size = Pt(10)
            run_title.font.color.rgb = RGBColor(0, 120, 212)
            
            run_text = p.add_run(text)
            run_text.font.name = 'Segoe UI'
            run_text.font.size = Pt(9)
            
    doc.add_paragraph().paragraph_format.space_after = Pt(10)
    
    # ------------------ SECTION 5 ------------------
    add_docx_heading(doc, "Section 5: Sample Client Journey", level=1)
    
    journey_table = doc.add_table(rows=1, cols=1)
    journey_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = journey_table.cell(0, 0)
    set_cell_bg(cell, "EFF6FF")
    set_callout_borders(cell, "107C41") # Green border for success case
    
    p = cell.paragraphs[0]
    p.paragraph_format.left_indent = Inches(0.15)
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(8)
    
    j_text = (
        "Integration Scenario:\n"
        "A new biopharma customer registers an assay CSV, PostgreSQL DB, Benchling ELN, and FASTA files.\n\n"
        "AnalytiX Automation:\n"
        "1. Establishes connectors to all data sources dynamically.\n"
        "2. Discovers database metadata attributes automatically (EAV catalog).\n"
        "3. Establishes chemical-biological relationships (SMILES mapping).\n"
        "4. Performs IC50 calculations and sequence alignments.\n"
        "5. Deploys dashboards, generates PDF reports, and enables the AI Scientist Chat interface."
    )
    j_run = p.add_run(j_text)
    j_run.font.name = 'Segoe UI'
    j_run.font.size = Pt(9.5)
    j_run.font.color.rgb = RGBColor(30, 58, 138)
    
    doc.add_paragraph().paragraph_format.space_after = Pt(10)
    
    # ------------------ SECTION 6 ------------------
    add_docx_heading(doc, "Section 6: Visual Value Chain", level=1)
    
    info_table = doc.add_table(rows=2, cols=4)
    info_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    set_table_light_borders(info_table)
    
    headers_info = ["CUSTOMER DATA", "PLATFORM INGEST", "AI & ANALYTICS", "OUTCOMES"]
    bg_colors = ["0F172A", "005A9E", "0078D4", "107C41"]
    
    for i, (h, bg) in enumerate(zip(headers_info, bg_colors)):
        cell = info_table.cell(0, i)
        set_cell_bg(cell, bg)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(5)
        p.paragraph_format.space_after = Pt(5)
        run = p.add_run(h)
        run.font.name = 'Segoe UI'
        run.font.bold = True
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(255, 255, 255)
        
    contents_info = [
        "ELN & LIMS databases\nCSV, Excel, FASTA\nCloud drop zones\nREST API payloads",
        "Connector discovery\nData validation\nMetadata Catalog\nFederated Engine",
        "AI Scientist Copilot\n4PL regressions\nAlignments grids\nCompound similarity",
        "Fast research loops\nRegulatory readiness\nIntegrated silos\nValidated inputs"
    ]
    
    for i, text in enumerate(contents_info):
        cell = info_table.cell(1, i)
        set_cell_bg(cell, "F8FAFC")
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(6)
        p.paragraph_format.space_after = Pt(6)
        run = p.add_run(text)
        run.font.name = 'Segoe UI'
        run.font.size = Pt(8.5)
        
    doc.save(filepath)
    print("DOCX generated successfully.")

# ---------------------------------------------------------
# PPTX GENERATION USING PYTHON-PPTX
# ---------------------------------------------------------

def set_shape_flat_style(shape, fill_color, line_color=None, line_width=1.5):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = PptxPt(line_width)
    else:
        shape.line.fill.background()

def add_pptx_textbox(slide, left, top, width, height, text, font_size=12, bold=False, color=PptxRGBColor(15, 23, 42), alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = PptxInches(0.02)
    tf.margin_bottom = PptxInches(0.02)
    tf.margin_left = PptxInches(0.02)
    tf.margin_right = PptxInches(0.02)
    
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = PptxPt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_pptx_bullets(slide, left, top, width, height, bullets, font_size=11, color=PptxRGBColor(51, 65, 85)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = PptxInches(0.02)
    tf.margin_bottom = PptxInches(0.02)
    tf.margin_left = PptxInches(0.02)
    tf.margin_right = PptxInches(0.02)
    
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_after = PptxPt(3)
        run = p.add_run()
        run.text = b
        run.font.name = 'Segoe UI'
        run.font.size = PptxPt(font_size)
        run.font.color.rgb = color
    return txBox

def create_pptx_slide_layout(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background full-screen rectangle (Light Slate gray)
    bg_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_flat_style(bg_rect, PptxRGBColor(248, 250, 252))
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptxInches(1.1))
    set_shape_flat_style(header, PptxRGBColor(11, 37, 64))
    
    # Title text
    add_pptx_textbox(slide, PptxInches(0.6), PptxInches(0.2), PptxInches(12.0), PptxInches(0.7), title_text, font_size=24, bold=True, color=PptxRGBColor(255, 255, 255))
    
    # Footer
    add_pptx_textbox(slide, PptxInches(0.6), PptxInches(7.1), PptxInches(8.0), PptxInches(0.3), "AnalytiX Platform - Inputs & Outputs Reference Deck", font_size=8, color=PptxRGBColor(100, 116, 139))
    
    return slide

def generate_pptx_document(filepath):
    print(f"Generating PPTX: {filepath}")
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    
    # ------------------ SLIDE 1: COVER ------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_flat_style(bg1, PptxRGBColor(11, 37, 64)) # Dark blue cover
    
    # Add title and subtitle
    add_pptx_textbox(slide1, PptxInches(0.8), PptxInches(1.8), PptxInches(11.5), PptxInches(1.5), "ANALYTIX PLATFORM", font_size=40, bold=True, color=PptxRGBColor(255, 255, 255))
    add_pptx_textbox(slide1, PptxInches(0.8), PptxInches(2.8), PptxInches(11.5), PptxInches(1.0), "Enterprise Inputs & Outputs Reference Guide", font_size=20, color=PptxRGBColor(0, 183, 195)) # Teal cyan
    
    # Accent line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(0.8), PptxInches(3.8), PptxInches(4.5), PptxInches(0.06))
    set_shape_flat_style(line, PptxRGBColor(0, 120, 212)) # Azure blue line
    
    # Metadata text box
    meta_text = (
        "Microsoft Azure Marketplace Technical Evaluation\n"
        "FDA 21 CFR Part 11 Compliant Environment\n"
        "Product Version: v4.2.0 (Production Release)"
    )
    add_pptx_textbox(slide1, PptxInches(0.8), PptxInches(5.0), PptxInches(8.0), PptxInches(1.5), meta_text, font_size=10.5, color=PptxRGBColor(148, 163, 184))
    
    # ------------------ SLIDE 2: INPUTS ------------------
    slide2 = create_pptx_slide_layout(prs, "Section 1: Supported Data Inputs")
    
    # 4 grid cards for categories
    cards_data = [
        ("Laboratory Systems", ["ELN: Benchling, Dotmatics", "LIMS: LabWare, LabVantage", "Assay definitions & audit records"], PptxInches(0.6), PptxInches(1.5)),
        ("Enterprise Databases", ["PostgreSQL, MySQL", "Oracle Database", "Microsoft SQL Server", "Snowflake Data Warehouse"], PptxInches(6.9), PptxInches(1.5)),
        ("Scientific Ingest Files", ["Structured: CSV, Excel", "Semi-structured: JSON, XML", "Biological: FASTA sequence files"], PptxInches(0.6), PptxInches(4.3)),
        ("APIs & Cloud Storage", ["REST & SOAP APIs integrations", "AWS S3, Azure Blob, GCS", "Secure SFTP drop zones"], PptxInches(6.9), PptxInches(4.3))
    ]
    
    for title, bullets, left, top in cards_data:
        # Background card shape
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, PptxInches(5.8), PptxInches(2.5))
        set_shape_flat_style(card, PptxRGBColor(255, 255, 255), PptxRGBColor(203, 213, 225))
        # Top banner on card
        card_header = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, PptxInches(5.8), PptxInches(0.5))
        set_shape_flat_style(card_header, PptxRGBColor(0, 120, 212))
        # Card title
        add_pptx_textbox(slide2, left + PptxInches(0.15), top + PptxInches(0.08), PptxInches(5.5), PptxInches(0.35), title, font_size=13, bold=True, color=PptxRGBColor(255, 255, 255))
        # Card bullets
        add_pptx_bullets(slide2, left + PptxInches(0.2), top + PptxInches(0.65), PptxInches(5.4), PptxInches(1.7), bullets, font_size=11)
        
    # ------------------ SLIDE 3: PROCESSING ------------------
    slide3 = create_pptx_slide_layout(prs, "Section 2: Platform Data Processing")
    
    # Ingestion flow at the top
    flow_steps = ["Data Input", "Connectors", "EAV Catalog", "Query Engine", "Analytics", "AI Scientist", "Reports & DB"]
    flow_colors = [
        PptxRGBColor(226, 232, 240),
        PptxRGBColor(203, 213, 225),
        PptxRGBColor(148, 163, 184),
        PptxRGBColor(100, 116, 139),
        PptxRGBColor(59, 130, 246),
        PptxRGBColor(29, 78, 216),
        PptxRGBColor(16, 124, 65)
    ]
    
    flow_width = PptxInches(1.5)
    flow_height = PptxInches(0.8)
    flow_top = PptxInches(1.5)
    
    for idx, (step, col) in enumerate(zip(flow_steps, flow_colors)):
        x_pos = PptxInches(0.6) + idx * PptxInches(1.75)
        # Rounded box
        flow_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, flow_top, flow_width, flow_height)
        set_shape_flat_style(flow_box, col)
        
        # Add text
        text_color = PptxRGBColor(255, 255, 255) if idx >= 4 else PptxRGBColor(15, 23, 42)
        add_pptx_textbox(slide3, x_pos, flow_top + PptxInches(0.15), flow_width, PptxInches(0.5), step, font_size=10, bold=True, color=text_color, alignment=PP_ALIGN.CENTER)
        
        # Draw arrow (if not last)
        if idx < len(flow_steps) - 1:
            arrow_x = x_pos + flow_width + PptxInches(0.05)
            add_pptx_textbox(slide3, arrow_x, flow_top + PptxInches(0.2), PptxInches(0.2), PptxInches(0.4), "➔", font_size=14, bold=True, color=PptxRGBColor(100, 116, 139), alignment=PP_ALIGN.CENTER)
            
    # List of processing actions below flow
    col_width = PptxInches(3.8)
    col_height = PptxInches(4.0)
    col_top = PptxInches(2.7)
    
    cols_data = [
        ("Ingestion & Federation", ["• Data validation & types check", "• Data quality filtering runs", "• Metadata Discovery & cataloging", "• Schema Discovery mapper", "• Federated queries on multiple DBs"], PptxInches(0.6)),
        ("Scientific Compute & AI", ["• Data integration & consolidation", "• Molecular compound calculations", "• DNA/RNA sequence alignments", "• Dose-response curve fitting (4PL)", "• AI Scientist Copilot processing"], PptxInches(4.75)),
        ("Audit & Compliance", ["• Metadata EAV modeling indexing", "• Compliance validation engine", "• Cryptographic audit logs chaining", "• Cryptographic SHA-256 signatures", "• Immutable data lineage tree"], PptxInches(8.9))
    ]
    
    for title, bullets, left in cols_data:
        # Title of column
        add_pptx_textbox(slide3, left, col_top, col_width, PptxInches(0.4), title, font_size=14, bold=True, color=PptxRGBColor(0, 90, 158))
        # Bullets
        add_pptx_bullets(slide3, left, col_top + PptxInches(0.45), col_width, col_height, bullets, font_size=11)
        
    # ------------------ SLIDE 4: OUTPUTS ------------------
    slide4 = create_pptx_slide_layout(prs, "Section 3: High-Value Platform Outputs")
    
    # 5 vertical cards
    outputs_cards = [
        ("Analytics Output", ["• Curve fit models (IC50)", "• Similarity matrices", "• Sequence identity plots", "• Scatter & Heatmap plots", "• Executive Dashboards"], PptxInches(0.6), PptxRGBColor(0, 120, 212)),
        ("AI Copilot", ["• Natural language answers", "• Research context brief", "• Molecular suggestions", "• Automated workflows", "• Executive Summaries"], PptxInches(3.0), PptxRGBColor(0, 90, 158)),
        ("Reports Ingest", ["• FDA compliance PDFs", "• Structured Excel runs", "• PowerPoint slide decks", "• CSV spreadsheet outputs", "• Regulatory audit packs"], PptxInches(5.4), PptxRGBColor(0, 183, 195)),
        ("Regulatory Logs", ["• Double-entry E-Signatures", "• SHA-256 audit ledger", "• Data lineage traces", "• Historical versions map", "• 21 CFR Part 11 summaries"], PptxInches(7.8), PptxRGBColor(118, 38, 143)),
        ("Operational Alerts", ["• Pipeline run updates", "• Approval requests", "• Custom Slack Webhooks", "• Immediate email alerts", "• Scheduled reports"], PptxInches(10.2), PptxRGBColor(16, 124, 65))
    ]
    
    for title, bullets, left, header_col in outputs_cards:
        # Card shape
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, PptxInches(1.5), PptxInches(2.1), PptxInches(5.3))
        set_shape_flat_style(card, PptxRGBColor(255, 255, 255), PptxRGBColor(226, 232, 240))
        # Header shape on card
        header_shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, PptxInches(1.5), PptxInches(2.1), PptxInches(0.6))
        set_shape_flat_style(header_shape, header_col)
        # Title text
        add_pptx_textbox(slide4, left, PptxInches(1.6), PptxInches(2.1), PptxInches(0.4), title, font_size=11, bold=True, color=PptxRGBColor(255, 255, 255), alignment=PP_ALIGN.CENTER)
        # Bullets
        add_pptx_bullets(slide4, left + PptxInches(0.08), PptxInches(2.2), PptxInches(1.94), PptxInches(4.4), bullets, font_size=10.5)
        
    # ------------------ SLIDE 5: BENEFITS ------------------
    slide5 = create_pptx_slide_layout(prs, "Section 4: Enterprise Business Benefits")
    
    benefits_cards = [
        ("Scientists", ["• Accelerates research cycles", "• Automates data aggregation", "• Interactive AI Copilot assistant"], PptxInches(0.6), PptxRGBColor(15, 118, 110)),
        ("Research Managers", ["• Unified program visibility", "• Real-time project dashboard", "• Scientific decision support"], PptxInches(3.7), PptxRGBColor(0, 120, 212)),
        ("Compliance Teams", ["• Full raw-to-final lineage", "• Cryptographic SHA-256 logs", "• Rapid FDA Part 11 reporting"], PptxInches(6.8), PptxRGBColor(118, 38, 143)),
        ("IT & Infrastructure", ["• Federated database query layers", "• Standard OAuth & RBAC rules", "• Scalable secure architecture"], PptxInches(9.9), PptxRGBColor(27, 38, 59))
    ]
    
    for title, bullets, left, header_col in benefits_cards:
        # Card shape
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, PptxInches(1.6), PptxInches(2.8), PptxInches(5.1))
        set_shape_flat_style(card, PptxRGBColor(255, 255, 255), PptxRGBColor(226, 232, 240))
        # Header shape on card
        header_shape = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, PptxInches(1.6), PptxInches(2.8), PptxInches(0.6))
        set_shape_flat_style(header_shape, header_col)
        # Title text
        add_pptx_textbox(slide5, left, PptxInches(1.7), PptxInches(2.8), PptxInches(0.4), title, font_size=13, bold=True, color=PptxRGBColor(255, 255, 255), alignment=PP_ALIGN.CENTER)
        # Bullets
        add_pptx_bullets(slide5, left + PptxInches(0.15), PptxInches(2.3), PptxInches(2.5), PptxInches(4.2), bullets, font_size=11)
        
    # ------------------ SLIDE 6: JOURNEY ------------------
    slide6 = create_pptx_slide_layout(prs, "Section 5: Sample Client Journey")
    
    # Left column: Customer registers
    left_x = PptxInches(0.6)
    top_y = PptxInches(1.6)
    width_w = PptxInches(5.8)
    height_h = PptxInches(5.1)
    
    card_left = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x, top_y, width_w, height_h)
    set_shape_flat_style(card_left, PptxRGBColor(254, 243, 199), PptxRGBColor(245, 158, 11)) # Amber light card
    add_pptx_textbox(slide6, left_x + PptxInches(0.3), top_y + PptxInches(0.3), width_w - PptxInches(0.6), PptxInches(0.5), "CUSTOMER DATA REGISTRATION", font_size=15, bold=True, color=PptxRGBColor(146, 64, 14))
    
    cust_bullets = [
        "1. Snowflake / PostgreSQL databases are connected via secure credentials.",
        "2. Benchling ELN research studies and assay run templates are imported.",
        "3. High-throughput CSV plates files containing screening numbers are uploaded.",
        "4. Genomic sequence FASTA files are added to sequence database registry."
    ]
    add_pptx_bullets(slide6, left_x + PptxInches(0.3), top_y + PptxInches(1.0), width_w - PptxInches(0.6), height_h - PptxInches(1.3), cust_bullets, font_size=12, color=PptxRGBColor(120, 83, 4))
    
    # Right column: AnalytiX automates
    right_x = PptxInches(6.9)
    card_right = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, top_y, width_w, height_h)
    set_shape_flat_style(card_right, PptxRGBColor(209, 250, 229), PptxRGBColor(16, 185, 129)) # Emerald light card
    add_pptx_textbox(slide6, right_x + PptxInches(0.3), top_y + PptxInches(0.3), width_w - PptxInches(0.6), PptxInches(0.5), "ANALYTIX PLATFORM AUTOMATION", font_size=15, bold=True, color=PptxRGBColor(6, 95, 70))
    
    plat_bullets = [
        "• Establish Read-Only Secure connections to all sources.",
        "• Automatically discover metadata structures (EAV Schema mapping).",
        "• Creates biological-chemical relationships (SMILES mappings).",
        "• Executes dose-response fits (IC50 curves fits) automatically.",
        "• Triggers analytics workflows and registers compound descriptors.",
        "• Deploys project dashboards, logs transactions, and launches AI Copilot."
    ]
    add_pptx_bullets(slide6, right_x + PptxInches(0.3), top_y + PptxInches(1.0), width_w - PptxInches(0.6), height_h - PptxInches(1.3), plat_bullets, font_size=12, color=PptxRGBColor(6, 78, 59))
    
    # ------------------ SLIDE 7: INFOGRAPHIC ------------------
    slide7 = create_pptx_slide_layout(prs, "Section 6: Visual Value Infographic")
    
    # 4 horizontal panels
    panel_width = PptxInches(2.7)
    panel_height = PptxInches(4.5)
    panel_top = PptxInches(1.8)
    
    panels = [
        ("CUSTOMER DATA", ["ELN & LIMS", "Databases & APIs", "Standard Files", "Raw SMILES structures"], PptxInches(0.6), PptxRGBColor(241, 245, 249), PptxRGBColor(15, 23, 42)),
        ("AnalytiX PLATFORM", ["Schema discovery", "Data validation", "EAV Catalog engine", "Data Lineage tracing"], PptxInches(3.7), PptxRGBColor(219, 234, 254), PptxRGBColor(30, 64, 175)),
        ("AI & ANALYTICS", ["AI Scientist Copilot", "Dose-Response fits", "Alignments grids", "Compound descriptors"], PptxInches(6.8), PptxRGBColor(30, 41, 59), PptxRGBColor(255, 255, 255)),
        ("BUSINESS VALUES", ["Accelerated Research", "100% FDA compliance", "No database silos", "Informed Decisions"], PptxInches(9.9), PptxRGBColor(6, 78, 59), PptxRGBColor(255, 255, 255))
    ]
    
    for title, bullets, left, bg_col, text_col in panels:
        # Card rounded rect
        panel_shape = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, panel_top, panel_width, panel_height)
        set_shape_flat_style(panel_shape, bg_col)
        
        # Title
        title_color = text_col if text_col != PptxRGBColor(255, 255, 255) else PptxRGBColor(0, 183, 195)
        add_pptx_textbox(slide7, left + PptxInches(0.15), panel_top + PptxInches(0.3), panel_width - PptxInches(0.3), PptxInches(0.5), title, font_size=13, bold=True, color=title_color, alignment=PP_ALIGN.CENTER)
        
        # Bullet list
        add_pptx_bullets(slide7, left + PptxInches(0.15), panel_top + PptxInches(1.0), panel_width - PptxInches(0.3), panel_height - PptxInches(1.2), bullets, font_size=11, color=text_col)
        
        # Add connecting arrows
        if left < PptxInches(9.0):
            arrow_left = left + panel_width + PptxInches(0.05)
            add_pptx_textbox(slide7, arrow_left, panel_top + PptxInches(2.0), PptxInches(0.3), PptxInches(0.5), "➔", font_size=16, bold=True, color=PptxRGBColor(100, 116, 139), alignment=PP_ALIGN.CENTER)
            
    prs.save(filepath)
    print("PowerPoint generated successfully.")

# ---------------------------------------------------------
# EXECUTION ENTRY POINT
# ---------------------------------------------------------

if __name__ == "__main__":
    pdf_path = os.path.join(DOCS_DIR, "AnalytiX_Platform_Inputs_Outputs.pdf")
    docx_path = os.path.join(DOCS_DIR, "AnalytiX_Platform_Inputs_Outputs.docx")
    pptx_path = os.path.join(DOCS_DIR, "AnalytiX_Platform_Inputs_Outputs.pptx")
    
    # Generate PDF
    generate_pdf_document(pdf_path)
    
    # Generate Word Document
    generate_docx_document(docx_path)
    
    # Generate PowerPoint Presentation
    generate_pptx_document(pptx_path)
    
    print("\nAll assets compiled successfully!")
    print(f"PDF: {pdf_path}")
    print(f"DOCX: {docx_path}")
    print(f"PowerPoint: {pptx_path}")
